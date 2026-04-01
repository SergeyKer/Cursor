import sys
import subprocess

try:
    from pptx import Presentation
except ImportError:
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx'])
    from pptx import Presentation

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Шпаргалка к собеседованию: GWARD"
slide.placeholders[1].text = "Подготовка к встречам с HR и Руководством\n(Конфиденциально, для кандидата)"

# Slide 2: Досье
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "1. Досье на компанию GWARD"
tf = slide.placeholders[1].text_frame
tf.text = "Ключевые вводные:"
p = tf.add_paragraph(); p.text = "- Масштаб: Производитель СИЗ рук (>140 моделей), отгружает >3 млн пар/мес."
p = tf.add_paragraph(); p.text = "- Лица: Максим Романенко (Владелец), Александр Ермаченко (Ген. директор)."
p = tf.add_paragraph(); p.text = "- Финансы: Выручка ~1.1-1.5 млрд руб. В 2024 году был спад выручки и минус по прибыли."
p = tf.add_paragraph(); p.text = "- Модель: Огромные запасы на складах (5-300 тыс. пар каждой модели) для отгрузки за 48 часов."

# Slide 3: Истинная боль (Инсайд)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "2. Инсайд: Почему они ищут Комдира"
tf = slide.placeholders[1].text_frame
tf.text = "Они уперлись в потолок старой модели продаж:"
p = tf.add_paragraph(); p.text = "- Замороженные деньги: Переизбыток товара на складах (ради отгрузки за 48ч) съедает P&L."
p = tf.add_paragraph(); p.text = "- Демпинг и дилеры: Отсутствие жесткой коммерческой политики и контроля РРЦ."
p = tf.add_paragraph(); p.text = "- Отдел продаж (30+ человек): Продают объем, а не маржу. Не хватает оцифровки и CRM."
p = tf.add_paragraph(); p.text = "--> ИМ НУЖЕН: Архитектор, который остановит 'разбазаривание' и выведет прибыль в плюс."

# Slide 4: Встреча с HR
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "3. Стратегия: Разговор с HR"
tf = slide.placeholders[1].text_frame
tf.text = "Цель: Показать адекватность, 'софты' и масштаб."
p = tf.add_paragraph(); p.text = "- Внутренний кандидат: Согласитесь, что им нужен человек с 'внешней насмотренностью'. Ваш опыт в Ansell и Portwest идеален."
p = tf.add_paragraph(); p.text = "- Позиционирование: Вы не 'супер-продавец', вы — системный управленец, который внедряет аналитику и CRM."
p = tf.add_paragraph(); p.text = "- Продукт: Знаете рынок СИЗ, понимаете, как забрать долю у ушедших 'европейцев'."

# Slide 5: Отработка 'Нюанса'
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "4. Легенда: Уход из Blesk InCare"
tf = slide.placeholders[1].text_frame
tf.text = "Как правильно объяснить ситуацию без негатива:"
p = tf.add_paragraph(); p.text = "1. '3 года успешно строил коммерцию (Текстиль), перевыполнял планы.'"
p = tf.add_paragraph(); p.text = "2. 'Доверили сложный кризисный проект по логистике перед самым сезоном с внедрением сырой 1С.'"
p = tf.add_paragraph(); p.text = "3. 'Вскрылся минус матричной структуры: ответственность на мне, а управление исполнителями (водителями) - у других.'"
p = tf.add_paragraph(); p.text = "4. 'Нам требовалась жесткая централизация, но руководство было не готово. Разошлись обоюдно. Моя страсть — Коммерция (B2B), поэтому я здесь.'"

# Slide 6: Встреча с ТОПами
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "5. Стратегия: Разговор с Топ-менеджерами"
tf = slide.placeholders[1].text_frame
tf.text = "Цель: Говорить на языке P&L (Прибыль/Убытки)."
p = tf.add_paragraph(); p.text = "- Аудит (1-й месяц): АВС/XYZ анализ 140 моделей ассортимента и всех дистрибьюторов."
p = tf.add_paragraph(); p.text = "- Маржа, а не оборот: Перевод мотивации РОПов на валовую прибыль и сбор дебиторки."
p = tf.add_paragraph(); p.text = "- Дилерская сеть: Жесткие грейды (скидки за объем и финансовую дисциплину)."
p = tf.add_paragraph(); p.text = "- Оборотный капитал: Переход от 'всем за 48 часов' к предзаказам для неключевых партнеров."

# Slide 7: Топ-вопросы
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "6. Ваши вопросы к HR (Киллер-фичи)"
tf = slide.placeholders[1].text_frame
tf.text = "Вопросы, чтобы показать системный подход:"
p = tf.add_paragraph(); p.text = "1. 'Как вы оцениваете корпоративную культуру? Она больше про жесткую иерархию или про гибкость и инициативу?'"
p = tf.add_paragraph(); p.text = "2. 'Какие 2-3 ключевые задачи будут стоять перед новым Комдиром на испытательный срок (первые 3 месяца)?'"
p = tf.add_paragraph(); p.text = "3. 'Есть ли уже утвержденный бюджет на автоматизацию (CRM/ERP), или мне предстоит его защищать с нуля?'"
p = tf.add_paragraph(); p.text = "4. 'Насколько текущая команда отдела продаж (30 чел) готова к изменениям и оцифровке процессов?'"

# Slide 8: Вопросы к Топам
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "7. Ваши вопросы к ТОПам (Владельцу/ГД)"
tf = slide.placeholders[1].text_frame
tf.text = "Спрашивайте в конце интервью с руководством:"
p = tf.add_paragraph(); p.text = "1. 'За счет чего вы видите кратный рост (x3)? Новые продукты, регионы или замена СТМ?'"
p = tf.add_paragraph(); p.text = "2. 'Как вы защищаете своих классических дилеров от ценового демпинга на маркетплейсах?'"
p = tf.add_paragraph(); p.text = "3. 'Готовы ли вы ради чистой прибыли отказаться от 'отгрузки за 48ч' для слабых клиентов?'"
p = tf.add_paragraph(); p.text = "4. 'Будет ли у меня карт-бланш на полное изменение системы KPI и, при необходимости, команды?'"

prs.save("GWARD_Candidate_CheatSheet_v2.pptx")
print("Presentation saved successfully.")
